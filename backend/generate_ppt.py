"""
CAPEX Portal - User Manual PowerPoint Generator
Creates a comprehensive PPT covering all roles and system workflows
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Color palette ──
PRIMARY = RGBColor(6, 182, 212)      # Cyan
SECONDARY = RGBColor(99, 102, 241)   # Indigo
DARK = RGBColor(15, 23, 42)          # Slate-900
DARK2 = RGBColor(30, 41, 59)         # Slate-800
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(241, 245, 249)      # Slate-100
MUTED = RGBColor(148, 163, 184)      # Slate-400
SUCCESS = RGBColor(16, 185, 129)     # Emerald
WARNING = RGBColor(245, 158, 11)     # Amber
ERROR = RGBColor(239, 68, 68)        # Red
ACCENT = RGBColor(139, 92, 246)      # Violet


def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=WHITE, spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(font_size * spacing)
        p.level = 0
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=12, font_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.name = 'Calibri'
    return shape


def add_flow_arrow(slide, left, top, width=0.6, height=0.3):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    return shape


def add_divider(slide, top, color=PRIMARY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(11.73), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_rounded_rect(slide, 0, 0, 13.33, 2.5, DARK2)
add_text_box(slide, 0.8, 0.8, 11.73, 1.0, "CAPEX PROCUREMENT PORTAL", 42, PRIMARY, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 1.6, 11.73, 0.6, "Complete User Manual & System Guide", 22, MUTED, False, PP_ALIGN.CENTER)
add_divider(slide, 2.5)
add_text_box(slide, 0.8, 3.2, 11.73, 0.5, "Streamline your capital expenditure lifecycle", 20, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 3.8, 11.73, 0.4, "From Request Submission to Commissioning", 16, MUTED, False, PP_ALIGN.CENTER)

# Feature boxes
features = [
    ("Role-Based Access", SUCCESS), ("Multi-Level Approvals", PRIMARY),
    ("Sample Management", ACCENT), ("AI Assistant", WARNING),
    ("Analytics Dashboard", SECONDARY), ("Master Data Mgmt", ERROR)
]
for i, (feat, col) in enumerate(features):
    x = 1.2 + (i % 3) * 3.8
    y = 4.8 + (i // 3) * 0.7
    add_rounded_rect(slide, x, y, 3.4, 0.5, col, feat, 12, WHITE, True)

add_text_box(slide, 0.8, 6.6, 11.73, 0.4, "Designed & Developed by Saurabh Jangir", 13, MUTED, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════
# SLIDE 2: TABLE OF CONTENTS
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11.73, 0.8, "TABLE OF CONTENTS", 32, PRIMARY, True)
add_divider(slide, 1.2)

toc_items = [
    "1.  System Overview & Architecture",
    "2.  Login & Authentication",
    "3.  User Roles & Permissions",
    "4.  Admin Panel - Master Data Management",
    "5.  Creating a CAPEX Request",
    "6.  Approval Workflow (Department Head)",
    "7.  Capex Head Review & Buyer Assignment",
    "8.  Buyer Workflow - Quotations & Purchase Orders",
    "9.  Sample Request & Dispatch Management",
    "10. Dashboard & Analytics",
    "11. Settings & Customization",
    "12. AI Assistant (Capex Man)",
    "13. Notifications & Communication",
]
add_bullet_list(slide, 1.0, 1.6, 11.0, 5.5, toc_items, 16, WHITE, 1.5)

# ════════════════════════════════════════════════
# SLIDE 3: SYSTEM OVERVIEW
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 8, 0.8, "1. SYSTEM OVERVIEW", 32, PRIMARY, True)
add_divider(slide, 1.2)
add_text_box(slide, 0.8, 1.5, 6, 0.4, "What is the CAPEX Portal?", 20, WHITE, True)
add_text_box(slide, 0.8, 2.0, 7, 1.2,
    "The CAPEX Procurement Portal is a comprehensive digital platform designed to manage "
    "the entire lifecycle of Capital Expenditure requests in manufacturing organizations. "
    "It replaces paper-based workflows with a streamlined, role-based digital process.",
    14, MUTED)

add_text_box(slide, 0.8, 3.5, 6, 0.4, "Key Capabilities", 18, SUCCESS, True)
caps = [
    "Digital CAPEX request submission with auto-populated fields",
    "Multi-level approval workflows (DH > Capex Head > Buyer)",
    "Supplier quotation management and comparison",
    "Purchase Order generation and tracking",
    "Sample request, dispatch, and inspection management",
    "Real-time analytics and cost savings tracking",
    "AI-powered assistant for instant query resolution",
    "Theme customization and dashboard personalization",
]
add_bullet_list(slide, 1.0, 3.9, 6.5, 3.0, caps, 13, LIGHT, 1.1)

# Tech stack box on right
add_rounded_rect(slide, 8.2, 1.5, 4.3, 5.0, DARK2)
add_text_box(slide, 8.4, 1.7, 4, 0.4, "Technology Stack", 16, PRIMARY, True)
stack = [
    "Frontend: React + Tailwind CSS",
    "Backend: FastAPI (Python)",
    "Database: MongoDB",
    "AI: OpenAI GPT-5 Integration",
    "Auth: JWT + Google OAuth",
    "UI: Shadcn/UI Components",
]
add_bullet_list(slide, 8.6, 2.3, 3.8, 2.5, stack, 13, MUTED, 1.5)

# ════════════════════════════════════════════════
# SLIDE 4: LOGIN & AUTHENTICATION
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 8, 0.8, "2. LOGIN & AUTHENTICATION", 32, PRIMARY, True)
add_divider(slide, 1.2)

add_text_box(slide, 0.8, 1.5, 6, 0.4, "How to Login", 20, WHITE, True)
login_steps = [
    "1. Open the portal URL in your browser",
    "2. The space-themed login page loads with an interactive bearing",
    "3. Click the pull cord OR press Space/Enter on keyboard",
    "4. The lamp lights up and the login form slides in from the right",
    "5. Enter your Email ID and Password",
    "6. Click 'Login' to access the portal",
    "   OR click 'Sign in with Google' for Google OAuth login",
]
add_bullet_list(slide, 1.0, 2.0, 5.5, 3.5, login_steps, 14, LIGHT, 1.3)

add_text_box(slide, 0.8, 5.2, 6, 0.4, "Forgot Password?", 16, WARNING, True)
forgot = [
    "Click 'Forgot password?' link on the login form",
    "Enter your registered email address",
    "Submit a reset request - Admin will review and approve",
    "Once approved, your password will be reset to default",
]
add_bullet_list(slide, 1.0, 5.7, 5.5, 1.5, forgot, 12, MUTED, 1.2)

# Right side info box
add_rounded_rect(slide, 7.5, 1.5, 5, 2.5, DARK2)
add_text_box(slide, 7.7, 1.7, 4.5, 0.4, "Login Methods", 16, SUCCESS, True)
add_bullet_list(slide, 7.9, 2.2, 4.3, 1.8, [
    "Email + Password (Standard)",
    "Google OAuth (Single Sign-On)",
    "Keyboard shortcut: Space / Enter to reveal form"
], 13, LIGHT, 1.5)

add_rounded_rect(slide, 7.5, 4.3, 5, 2.8, DARK2)
add_text_box(slide, 7.7, 4.5, 4.5, 0.4, "Security Features", 16, ERROR, True)
add_bullet_list(slide, 7.9, 5.0, 4.3, 2.0, [
    "JWT token-based authentication",
    "Admin-controlled password resets",
    "Session management with auto-logout",
    "Brute force protection",
    "Role-based page access control"
], 13, MUTED, 1.3)

# ════════════════════════════════════════════════
# SLIDE 5: USER ROLES
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 8, 0.8, "3. USER ROLES & PERMISSIONS", 32, PRIMARY, True)
add_divider(slide, 1.2)

roles = [
    ("Admin", ERROR, [
        "Full system administration access",
        "Manage Plants, Departments, Users",
        "Approve/reject password reset requests",
        "Configure themes for admin panel",
        "View system-wide statistics",
    ]),
    ("Capex Head", PRIMARY, [
        "Review and approve/reject CAPEX requests",
        "Assign buyers to approved requests",
        "Monitor all requests and spending",
        "Access executive analytics dashboard",
        "Edit user details in sidebar",
    ]),
    ("Buyer", SUCCESS, [
        "Manage assigned CAPEX requests",
        "Handle supplier quotations and comparisons",
        "Generate and manage Purchase Orders",
        "Initiate and manage Sample Requests",
        "Track dispatch, gate pass, and delivery",
    ]),
]

for i, (role, color, perms) in enumerate(roles):
    x = 0.8 + i * 4.1
    add_rounded_rect(slide, x, 1.5, 3.8, 0.5, color, role, 16, WHITE, True)
    add_bullet_list(slide, x + 0.2, 2.2, 3.4, 2.5, perms, 11, LIGHT, 1.2)

roles2 = [
    ("Department Head", WARNING, [
        "Approve/reject requests from department users",
        "View department-specific analytics",
        "Monitor budget utilization",
        "Access detailed request information",
    ]),
    ("Process Engineer", ACCENT, [
        "Submit CAPEX requests for equipment/tooling",
        "Track request progress through stages",
        "View technical specifications and samples",
        "Provide engineering input on requirements",
    ]),
    ("User", MUTED, [
        "Submit new CAPEX requests",
        "Track personal request status",
        "View assigned request details",
        "Receive notifications on updates",
    ]),
]

for i, (role, color, perms) in enumerate(roles2):
    x = 0.8 + i * 4.1
    add_rounded_rect(slide, x, 4.7, 3.8, 0.5, color, role, 16, WHITE, True)
    add_bullet_list(slide, x + 0.2, 5.4, 3.4, 2.0, perms, 11, LIGHT, 1.2)

# ════════════════════════════════════════════════
# SLIDE 6: ADMIN PANEL
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 8, 0.8, "4. ADMIN PANEL", 32, PRIMARY, True)
add_text_box(slide, 0.8, 0.9, 8, 0.4, "Master Data Management", 18, MUTED)
add_divider(slide, 1.4)

# Plants section
add_rounded_rect(slide, 0.8, 1.7, 3.8, 2.8, DARK2)
add_text_box(slide, 1.0, 1.8, 3.4, 0.4, "Plants Management", 16, WARNING, True)
add_bullet_list(slide, 1.2, 2.3, 3.2, 2.0, [
    "Add new manufacturing plants",
    "Edit existing plant names",
    "Delete plants (if no dependencies)",
    "Plants are linked to departments",
], 12, LIGHT, 1.3)

# Departments section
add_rounded_rect(slide, 4.9, 1.7, 3.8, 2.8, DARK2)
add_text_box(slide, 5.1, 1.8, 3.4, 0.4, "Departments", 16, SECONDARY, True)
add_bullet_list(slide, 5.3, 2.3, 3.2, 2.0, [
    "Create departments under plants",
    "Assign departments to specific plants",
    "Edit department names and plant mapping",
    "Department-based request routing",
], 12, LIGHT, 1.3)

# Users section
add_rounded_rect(slide, 9.0, 1.7, 3.8, 2.8, DARK2)
add_text_box(slide, 9.2, 1.8, 3.4, 0.4, "User Management", 16, SUCCESS, True)
add_bullet_list(slide, 9.4, 2.3, 3.2, 2.0, [
    "Create user accounts with all details",
    "Assign roles, plant, department",
    "Map users to Department Heads",
    "Reset passwords and manage access",
], 12, LIGHT, 1.3)

add_text_box(slide, 0.8, 4.8, 6, 0.4, "Creating a New User", 18, WHITE, True)
user_fields = [
    "Required fields: Email, Full Name, Employee ID, Mobile Number",
    "Account Setup: Password (for new users), Role selection",
    "Location: Plant & Department assignment (for User / DH roles)",
    "Reporting: Map to Department Head (User role) or Capex Head (Buyer role)",
    "The form is organized into labeled sections for professional data entry",
]
add_bullet_list(slide, 1.0, 5.3, 7.0, 2.0, user_fields, 13, LIGHT, 1.2)

# Theme box
add_rounded_rect(slide, 8.5, 4.8, 4.3, 2.2, DARK2)
add_text_box(slide, 8.7, 4.9, 4, 0.4, "Admin Theme", 16, ACCENT, True)
add_bullet_list(slide, 8.9, 5.4, 3.8, 1.5, [
    "Choose from multiple themes",
    "Theme applies to admin panel only",
    "Visual preview before selection",
    "Saved per admin session",
], 12, MUTED, 1.3)

# ════════════════════════════════════════════════
# SLIDE 7: CREATING A CAPEX REQUEST
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 8, 0.8, "5. CREATING A CAPEX REQUEST", 32, PRIMARY, True)
add_divider(slide, 1.2)

add_text_box(slide, 0.8, 1.5, 6, 0.4, "Step-by-Step Request Submission", 20, WHITE, True)
steps = [
    "1. Click 'New Request' from the Dashboard or sidebar menu",
    "2. Plant and Department are auto-populated from your profile (non-editable)",
    "3. Fill in the request details:",
    "     - Item Name / Description / Specification",
    "     - Quantity and Estimated Budget",
    "     - Justification for the expenditure",
    "     - Priority level (Low / Medium / High / Critical)",
    "4. Attach supporting documents if needed (drawings, specs, photos)",
    "5. Submit the request for Department Head approval",
    "6. Track your request status from the Dashboard",
]
add_bullet_list(slide, 1.0, 2.0, 6.0, 4.5, steps, 14, LIGHT, 1.2)

# Side info
add_rounded_rect(slide, 7.8, 1.5, 4.7, 3.0, DARK2)
add_text_box(slide, 8.0, 1.7, 4.3, 0.4, "Auto-Populated Fields", 16, SUCCESS, True)
add_bullet_list(slide, 8.2, 2.2, 4.0, 2.0, [
    "Plant - from your user profile",
    "Department - from your user profile",
    "Requested By - your name",
    "Date - auto-set to current date",
    "Status - auto-set to 'Pending'",
], 13, MUTED, 1.3)

add_rounded_rect(slide, 7.8, 4.8, 4.7, 2.2, DARK2)
add_text_box(slide, 8.0, 5.0, 4.3, 0.4, "Request Statuses", 16, WARNING, True)
statuses = ["Pending", "DH Approved", "CH Approved", "Buyer Assigned", "Quoted", "PO Generated", "Completed"]
add_bullet_list(slide, 8.2, 5.5, 4.0, 1.5, statuses, 12, LIGHT, 1.1)

# ════════════════════════════════════════════════
# SLIDE 8: APPROVAL WORKFLOW
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 8, 0.8, "6. APPROVAL WORKFLOW", 32, PRIMARY, True)
add_divider(slide, 1.2)

add_text_box(slide, 0.8, 1.5, 12, 0.4, "Request Lifecycle Flow", 20, WHITE, True)

# Flow diagram
flow_items = [
    ("User\nSubmits", MUTED, 0.8),
    ("DH\nReview", WARNING, 2.8),
    ("Capex Head\nApproval", PRIMARY, 4.8),
    ("Buyer\nAssigned", SUCCESS, 6.8),
    ("Quotation\n& Comparison", ACCENT, 8.8),
    ("PO\nGenerated", SECONDARY, 10.8),
]
for (label, color, x) in flow_items:
    add_rounded_rect(slide, x, 2.2, 1.6, 0.9, color, label, 10, WHITE, True)
    if x < 10.8:
        add_flow_arrow(slide, x + 1.6, 2.5)

# DH section
add_rounded_rect(slide, 0.8, 3.6, 5.8, 3.2, DARK2)
add_text_box(slide, 1.0, 3.7, 5.4, 0.4, "Department Head Actions", 18, WARNING, True)
add_bullet_list(slide, 1.2, 4.2, 5.2, 2.5, [
    "View all pending requests from department users",
    "Open request detail to review specifications",
    "Approve: Request moves to Capex Head review",
    "Reject: Request sent back to user with comments",
    "Requests appear in 'Pending Tasks' on dashboard",
    "Can add remarks/comments before approving/rejecting",
], 13, LIGHT, 1.2)

# CH section
add_rounded_rect(slide, 7.0, 3.6, 5.8, 3.2, DARK2)
add_text_box(slide, 7.2, 3.7, 5.4, 0.4, "Capex Head Actions", 18, PRIMARY, True)
add_bullet_list(slide, 7.4, 4.2, 5.2, 2.5, [
    "Review DH-approved requests for final authorization",
    "Approve: Enables buyer assignment for procurement",
    "Reject: Request returned with rejection reason",
    "Assign Buyer: Select from available buyers",
    "Monitor all requests across departments",
    "Track cost savings and budget utilization",
], 13, LIGHT, 1.2)

# ════════════════════════════════════════════════
# SLIDE 9: BUYER WORKFLOW
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 8, 0.8, "8. BUYER WORKFLOW", 32, PRIMARY, True)
add_text_box(slide, 0.8, 0.9, 8, 0.4, "Quotations, Comparison & Purchase Orders", 18, MUTED)
add_divider(slide, 1.4)

# Quotation
add_rounded_rect(slide, 0.8, 1.7, 3.8, 3.0, DARK2)
add_text_box(slide, 1.0, 1.8, 3.4, 0.4, "1. Supplier Quotations", 16, SUCCESS, True)
add_bullet_list(slide, 1.2, 2.3, 3.2, 2.2, [
    "Add supplier details (name, contact)",
    "Enter quoted prices (initial & final)",
    "Upload quotation documents",
    "Add multiple suppliers for comparison",
    "Track negotiation progress",
], 12, LIGHT, 1.2)

# Comparison
add_rounded_rect(slide, 4.9, 1.7, 3.8, 3.0, DARK2)
add_text_box(slide, 5.1, 1.8, 3.4, 0.4, "2. Price Comparison", 16, WARNING, True)
add_bullet_list(slide, 5.3, 2.3, 3.2, 2.2, [
    "Side-by-side supplier comparison",
    "Compare initial vs final prices",
    "Calculate cost savings automatically",
    "Select preferred supplier",
    "Justify selection with comments",
], 12, LIGHT, 1.2)

# PO
add_rounded_rect(slide, 9.0, 1.7, 3.8, 3.0, DARK2)
add_text_box(slide, 9.2, 1.8, 3.4, 0.4, "3. Purchase Orders", 16, PRIMARY, True)
add_bullet_list(slide, 9.4, 2.3, 3.2, 2.2, [
    "Generate PO from approved quotation",
    "PO number auto-assigned",
    "Track PO status and delivery dates",
    "Upload PO documents",
    "Link PO to CAPEX request",
], 12, LIGHT, 1.2)

add_text_box(slide, 0.8, 5.0, 6, 0.4, "Buyer Dashboard Features", 18, WHITE, True)
add_bullet_list(slide, 1.0, 5.5, 11, 1.5, [
    "Pending Tasks widget shows all assigned requests requiring action",
    "Quick filters by status (Quoted, PO Pending, Sample Required, etc.)",
    "Buyer-specific analytics: quotations sent, POs generated, cost savings achieved",
    "Sample management: initiate, track dispatch, verify delivery",
], 13, LIGHT, 1.2)

# ════════════════════════════════════════════════
# SLIDE 10: SAMPLE MANAGEMENT
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 8, 0.8, "9. SAMPLE MANAGEMENT", 32, PRIMARY, True)
add_text_box(slide, 0.8, 0.9, 8, 0.4, "Request, Dispatch & Inspection Workflow", 18, MUTED)
add_divider(slide, 1.4)

# Flow
sample_flow = [
    ("Sample\nRequested", ACCENT, 0.5),
    ("Under\nPreparation", WARNING, 2.7),
    ("Ready for\nDispatch", PRIMARY, 4.9),
    ("Gate Pass\nIssued", SUCCESS, 7.1),
    ("Dispatched", SECONDARY, 9.3),
    ("Delivered", RGBColor(34, 197, 94), 11.3),
]
for (label, color, x) in sample_flow:
    add_rounded_rect(slide, x, 1.7, 1.7, 0.8, color, label, 10, WHITE, True)
    if x < 11.3:
        add_flow_arrow(slide, x + 1.7, 1.95, 0.5)

add_rounded_rect(slide, 0.8, 3.0, 5.8, 3.8, DARK2)
add_text_box(slide, 1.0, 3.1, 5.4, 0.4, "Sample Request Details", 16, ACCENT, True)
add_bullet_list(slide, 1.2, 3.6, 5.2, 3.0, [
    "Linked to CAPEX request and supplier",
    "Specify quantity, packing type, special instructions",
    "Track preparation status with buyer updates",
    "Upload challan and gate pass documents",
    "Record dispatch details (vehicle, driver, tracking)",
    "Delivery confirmation with date and photos",
    "Visibility: Buyers see only their own samples",
    "Others see read-only sample details via request view",
], 13, LIGHT, 1.1)

add_rounded_rect(slide, 7.0, 3.0, 5.8, 3.8, DARK2)
add_text_box(slide, 7.2, 3.1, 5.4, 0.4, "Sample Card Information", 16, PRIMARY, True)
add_bullet_list(slide, 7.4, 3.6, 5.2, 3.0, [
    "Sample ID and linked request number",
    "Supplier name and contact details",
    "Quantity and packing type",
    "Current status with color-coded badges",
    "Preparation and dispatch timestamps",
    "Gate pass number and challan reference",
    "File attachments (documents, photos)",
    "Cards support pagination for large lists",
], 13, LIGHT, 1.1)

# ════════════════════════════════════════════════
# SLIDE 11: DASHBOARD & ANALYTICS
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 8, 0.8, "10. DASHBOARD & ANALYTICS", 32, PRIMARY, True)
add_divider(slide, 1.2)

add_rounded_rect(slide, 0.8, 1.5, 3.6, 2.2, DARK2)
add_text_box(slide, 1.0, 1.6, 3.2, 0.4, "Statistics Cards", 16, SUCCESS, True)
add_bullet_list(slide, 1.2, 2.1, 3.0, 1.5, [
    "Total Requests count",
    "Pending approvals count",
    "Active POs in pipeline",
    "Cost savings achieved",
    "Role-specific metrics",
], 12, LIGHT, 1.1)

add_rounded_rect(slide, 4.7, 1.5, 3.8, 2.2, DARK2)
add_text_box(slide, 4.9, 1.6, 3.4, 0.4, "Executive Analytics", 16, WARNING, True)
add_bullet_list(slide, 5.1, 2.1, 3.2, 1.5, [
    "Department spend analysis (chart)",
    "Budget vs actual comparison",
    "Status breakdown pie chart",
    "Buyer performance metrics",
    "Cost savings by supplier",
], 12, LIGHT, 1.1)

add_rounded_rect(slide, 8.8, 1.5, 3.8, 2.2, DARK2)
add_text_box(slide, 9.0, 1.6, 3.4, 0.4, "Recent Requests", 16, PRIMARY, True)
add_bullet_list(slide, 9.2, 2.1, 3.2, 1.5, [
    "Sortable table with all requests",
    "Click to view full details",
    "Status badges with colors",
    "Quick search and filters",
    "Pagination for large datasets",
], 12, LIGHT, 1.1)

add_text_box(slide, 0.8, 4.0, 6, 0.4, "Dashboard Customization", 20, WHITE, True)
add_bullet_list(slide, 1.0, 4.5, 11, 2.5, [
    "Access Settings (gear icon in sidebar) > Dashboard tab",
    "Toggle widget visibility ON/OFF for a personalized dashboard",
    "Drag and reorder widgets to prioritize what matters most",
    "Available widgets: Statistics Cards, Executive Analytics, Pending Tasks, Requests Table",
    "Customizations are saved per user and persist across sessions",
    "Reset to default layout anytime with one click",
], 14, LIGHT, 1.3)

# ════════════════════════════════════════════════
# SLIDE 12: SETTINGS & CUSTOMIZATION
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 8, 0.8, "11. SETTINGS & CUSTOMIZATION", 32, PRIMARY, True)
add_divider(slide, 1.2)

# Appearance
add_rounded_rect(slide, 0.8, 1.5, 5.8, 2.5, DARK2)
add_text_box(slide, 1.0, 1.6, 5.4, 0.4, "Appearance Settings", 18, ACCENT, True)
add_bullet_list(slide, 1.2, 2.1, 5.2, 1.8, [
    "Choose from multiple color themes (Aurora, Sunset, Ocean, etc.)",
    "Theme changes apply to the entire user dashboard",
    "Adjust font size preferences (Small, Medium, Large)",
    "Visual preview of each theme with color swatches",
    "Changes saved instantly and persist across sessions",
], 13, LIGHT, 1.2)

# Security
add_rounded_rect(slide, 7.0, 1.5, 5.8, 2.5, DARK2)
add_text_box(slide, 7.2, 1.6, 5.4, 0.4, "Security Settings", 18, ERROR, True)
add_bullet_list(slide, 7.4, 2.1, 5.2, 1.8, [
    "Change password (current + new + confirm)",
    "Password strength requirements enforced",
    "Session management information",
    "Last login timestamp display",
    "Account security recommendations",
], 13, LIGHT, 1.2)

# Dashboard
add_rounded_rect(slide, 0.8, 4.3, 5.8, 2.8, DARK2)
add_text_box(slide, 1.0, 4.4, 5.4, 0.4, "Dashboard Settings", 18, SUCCESS, True)
add_bullet_list(slide, 1.2, 4.9, 5.2, 2.0, [
    "Toggle individual widget visibility",
    "Drag widgets up/down to reorder",
    "Available: Stats, Analytics, Pending Tasks, Requests Table",
    "Save layout configuration",
    "Reset to default with one click",
], 13, LIGHT, 1.2)

# Help
add_rounded_rect(slide, 7.0, 4.3, 5.8, 2.8, DARK2)
add_text_box(slide, 7.2, 4.4, 5.4, 0.4, "Help & Support", 18, WARNING, True)
add_bullet_list(slide, 7.4, 4.9, 5.2, 2.0, [
    "Start Interactive Tutorial (11-step walkthrough)",
    "Download User Manual (this PPT)",
    "Tooltips guide new users through the interface",
    "Context-sensitive help throughout the app",
    "AI Assistant available for instant answers",
], 13, LIGHT, 1.2)

# ════════════════════════════════════════════════
# SLIDE 13: AI ASSISTANT
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 8, 0.8, "12. AI ASSISTANT - CAPEX MAN", 32, PRIMARY, True)
add_divider(slide, 1.2)

add_text_box(slide, 0.8, 1.5, 6, 0.4, "Your Intelligent Portal Companion", 20, WHITE, True)
add_bullet_list(slide, 1.0, 2.0, 6, 3.5, [
    "Powered by OpenAI GPT-5 for intelligent responses",
    "Click the 'Capex Man' floating button on the dashboard",
    "Ask questions in natural language about:",
    "     - How to submit a request",
    "     - Status of your pending approvals",
    "     - Understanding approval workflows",
    "     - Cost savings and budget queries",
    "     - System navigation and feature guidance",
    "Context-aware: understands your role and permissions",
    "Conversation history maintained within session",
], 14, LIGHT, 1.2)

add_rounded_rect(slide, 7.5, 1.5, 5, 3.0, DARK2)
add_text_box(slide, 7.7, 1.7, 4.6, 0.4, "Example Queries", 16, SUCCESS, True)
add_bullet_list(slide, 7.9, 2.2, 4.3, 2.2, [
    '"How do I create a new CAPEX request?"',
    '"What is the status of my pending approvals?"',
    '"Explain the sample dispatch process"',
    '"How do I compare supplier quotations?"',
    '"What permissions does a Buyer have?"',
], 12, MUTED, 1.3)

add_rounded_rect(slide, 7.5, 4.8, 5, 2.0, DARK2)
add_text_box(slide, 7.7, 5.0, 4.6, 0.4, "AI Capabilities", 16, ACCENT, True)
add_bullet_list(slide, 7.9, 5.5, 4.3, 1.3, [
    "Natural language understanding",
    "Role-aware responses",
    "Process guidance and tutorials",
    "Data-driven answers from portal context",
], 12, LIGHT, 1.2)

# ════════════════════════════════════════════════
# SLIDE 14: NOTIFICATIONS
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 8, 0.8, "13. NOTIFICATIONS", 32, PRIMARY, True)
add_divider(slide, 1.2)

add_text_box(slide, 0.8, 1.5, 6, 0.4, "Stay Updated on Request Progress", 20, WHITE, True)

add_rounded_rect(slide, 0.8, 2.2, 5.8, 4.5, DARK2)
add_text_box(slide, 1.0, 2.3, 5.4, 0.4, "Notification Types", 18, SUCCESS, True)
add_bullet_list(slide, 1.2, 2.8, 5.2, 3.5, [
    "Request Submitted - confirmation to submitter",
    "Approval Required - alert to DH/Capex Head",
    "Request Approved/Rejected - status update",
    "Buyer Assigned - notification to buyer",
    "Quotation Added - update to stakeholders",
    "PO Generated - procurement notification",
    "Sample Status Change - dispatch updates",
    "Password Reset - admin notification",
    "System Announcements - portal-wide updates",
], 13, LIGHT, 1.2)

add_rounded_rect(slide, 7.0, 2.2, 5.8, 4.5, DARK2)
add_text_box(slide, 7.2, 2.3, 5.4, 0.4, "How Notifications Work", 18, PRIMARY, True)
add_bullet_list(slide, 7.4, 2.8, 5.2, 3.5, [
    "Bell icon in the top header shows unread count",
    "Click bell to see notification dropdown",
    "Unread notifications highlighted with blue background",
    "Click a notification to navigate to the relevant page",
    "Mark all as read with one click",
    "Notifications auto-generated for key workflow events",
    "Real-time updates without page refresh",
    "Up to 10 most recent notifications shown",
    "Older notifications accessible via notification page",
], 13, LIGHT, 1.2)

# ════════════════════════════════════════════════
# SLIDE 15: QUICK REFERENCE
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 8, 0.8, "QUICK REFERENCE CARD", 32, PRIMARY, True)
add_divider(slide, 1.2)

add_text_box(slide, 0.8, 1.5, 6, 0.4, "Navigation Shortcuts", 18, WHITE, True)

# Left column
add_rounded_rect(slide, 0.8, 2.0, 3.8, 4.8, DARK2)
add_text_box(slide, 1.0, 2.1, 3.4, 0.4, "Sidebar Menu", 14, SUCCESS, True)
add_bullet_list(slide, 1.2, 2.5, 3.2, 4.0, [
    "Dashboard - Home/overview",
    "All Requests - Browse all CAPEX requests",
    "Samples - Sample request management",
    "New Request - Submit new CAPEX",
    "Notifications - View all notifications",
    "User Details - Profile & user info",
    "Settings - Customize & configure",
    "",
    "Sidebar can be collapsed/expanded",
    "with the toggle button",
    "",
    "On mobile: Use hamburger menu",
], 11, LIGHT, 1.0)

# Middle column
add_rounded_rect(slide, 4.9, 2.0, 3.8, 4.8, DARK2)
add_text_box(slide, 5.1, 2.1, 3.4, 0.4, "Keyboard Shortcuts", 14, WARNING, True)
add_bullet_list(slide, 5.3, 2.5, 3.2, 4.0, [
    "Space / Enter - Open login form",
    "",
    "Common Actions:",
    "Click card - View request details",
    "Click status badge - See status info",
    "Click bell icon - View notifications",
    "Click profile - User dropdown menu",
    "",
    "Filters:",
    "Status dropdown - Filter by status",
    "Buyer dropdown - Filter by buyer",
    "Search bar - Quick text search",
], 11, LIGHT, 1.0)

# Right column
add_rounded_rect(slide, 9.0, 2.0, 3.8, 4.8, DARK2)
add_text_box(slide, 9.2, 2.1, 3.4, 0.4, "Mobile Support", 14, PRIMARY, True)
add_bullet_list(slide, 9.4, 2.5, 3.2, 4.0, [
    "Fully responsive design",
    "Works on all screen sizes",
    "",
    "Mobile Features:",
    "Hamburger menu for navigation",
    "Touch-friendly card interfaces",
    "Swipe-friendly tab navigation",
    "Stacked layouts on small screens",
    "Scrollable tables with touch",
    "",
    "Optimized for:",
    "iPhone, Android, Tablets, iPad",
], 11, LIGHT, 1.0)

# ════════════════════════════════════════════════
# SLIDE 16: THANK YOU
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rounded_rect(slide, 0, 0, 13.33, 7.5, DARK)
add_text_box(slide, 0.8, 2.0, 11.73, 1.0, "CAPEX PROCUREMENT PORTAL", 42, PRIMARY, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 3.0, 11.73, 0.6, "User Manual", 24, MUTED, False, PP_ALIGN.CENTER)
add_divider(slide, 3.8)
add_text_box(slide, 0.8, 4.2, 11.73, 0.5, "For support, contact your system administrator", 16, LIGHT, False, PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 4.8, 11.73, 0.5, "or use the AI Assistant (Capex Man) for instant help", 16, MUTED, False, PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 5.8, 11.73, 0.5, "Designed & Developed by Saurabh Jangir", 18, PRIMARY, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 6.4, 11.73, 0.4, "Version 1.0", 12, MUTED, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════
output_path = "/app/backend/static/CAPEX_Portal_User_Manual.pptx"
prs.save(output_path)
print(f"PPT saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
